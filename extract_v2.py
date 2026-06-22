
import sys, os, traceback
sys.path.insert(0, r'G:\ppt-master-main\skills\ppt-master\scripts')
from pptx import Presentation

pptx_path = r'G:\ppt-master-main\projects\GeCoM-Net_ppt169_20240616\exports\GeCoM-Net_修改版.pptx'
out_dir = r'G:\ppt-master-main\projects\GeCoM-Net_ppt169_20240616\validation\slide_images'
os.makedirs(out_dir, exist_ok=True)
out_file = r'G:\ppt-master-main\projects\GeCoM-Net_ppt169_20240616\validation\slide_analysis.txt'
lines = []

prs = Presentation(pptx_path)
slide12 = prs.slides[11]
lines.append('=== Slide 12 Images ===')
for shape in slide12.shapes:
    if shape.shape_type == 13:
        ext = shape.image.content_type.split('/')[-1]
        fname = 'slide12_' + shape.name + '.' + ext
        fpath = os.path.join(out_dir, fname)
        with open(fpath, 'wb') as f:
            f.write(shape.image.blob)
        lines.append('Saved: %s (%d bytes) pos=(%d,%d) size=(%d,%d)' % (fname, len(shape.image.blob), shape.left, shape.top, shape.width, shape.height))

lines.append('')
lines.append('=== Slide 9 detailed ===')
slide9 = prs.slides[8]

def dump(shape, indent=0):
    prefix = '  ' * indent
    lines.append('%s%d:%s t=%s pos=(%d,%d) sz=(%d,%d)' % (prefix, shape.shape_id, shape.name, shape.shape_type, shape.left, shape.top, shape.width, shape.height))
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            ri = []
            for r in p.runs:
                ri.append('[%s|b=%s,s=%s]' % (r.text, r.font.bold, r.font.size))
            if p.text.strip():
                lines.append('%s  P: %s' % (prefix, p.text[:150]))
                if ri:
                    lines.append('%s  R: %s' % (prefix, ' '.join(ri)[:300]))
    if hasattr(shape, 'shapes'):
        try:
            for c in shape.shapes:
                dump(c, indent+1)
        except:
            pass

for shape in slide9.shapes:
    dump(shape)

lines.append('')
lines.append('=== Slide 12 detailed ===')
for shape in slide12.shapes:
    dump(shape)

with open(out_file, 'w', encoding='utf-8') as f:
    f.write('\n'.join(lines))
print('OK')
