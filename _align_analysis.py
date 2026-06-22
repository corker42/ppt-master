#!/usr/bin/env python3
"""Analyze alignment of both target and reference PPTs in detail."""
import sys, os
sys.path.insert(0, 'skills/ppt-master/scripts')
from pptx import Presentation
from pptx.util import Emu, Pt

def emu_to_cm(emu):
    return round(emu / 360000, 2)

def get_all_shapes_recursive(shape, depth=0):
    """Get all leaf shapes with their global position."""
    results = []
    if shape.shape_type == 6:  # GROUP
        # For groups, we need to consider the group's position
        for child in shape.shapes:
            child_results = get_all_shapes_recursive(child, depth + 1)
            # Offset child positions by group position
            for cr in child_results:
                cr['global_left'] = shape.left + cr.get('global_left', cr['left'])
                cr['global_top'] = shape.top + cr.get('global_top', cr['top'])
            results.extend(child_results)
    else:
        results.append({
            'name': shape.name,
            'type': str(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height,
            'global_left': shape.left,
            'global_top': shape.top,
            'depth': depth,
            'text': shape.text_frame.text[:40] if shape.has_text_frame else '',
            'is_image': shape.shape_type == 13,
        })
    return results

def analyze_slide(prs, slide_idx, label):
    slide = prs.slides[slide_idx]
    canvas_w = prs.slide_width
    canvas_h = prs.slide_height
    
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_idx+1} ({label})")
    print(f"Canvas: {emu_to_cm(canvas_w)}cm x {emu_to_cm(canvas_h)}cm ({canvas_w} x {canvas_h} EMU)")
    print(f"{'='*60}")
    
    all_shapes = []
    for shape in slide.shapes:
        shapes = get_all_shapes_recursive(shape)
        for s in shapes:
            s['parent'] = shape.name
        all_shapes.extend(shapes)
    
    # Sort by position
    all_shapes.sort(key=lambda s: (s['global_top'], s['global_left']))
    
    for s in all_shapes:
        left_cm = emu_to_cm(s['global_left'])
        top_cm = emu_to_cm(s['global_top'])
        w_cm = emu_to_cm(s['width'])
        h_cm = emu_to_cm(s['height'])
        right_cm = emu_to_cm(s['global_left'] + s['width'])
        bottom_cm = emu_to_cm(s['global_top'] + s['height'])
        
        img_flag = ' [IMG]' if s['is_image'] else ''
        text_preview = f' "{s["text"]}"' if s['text'] else ''
        
        print(f"  {s['name']:25s} L={left_cm:6.2f} T={top_cm:6.2f} W={w_cm:6.2f} H={h_cm:6.2f} R={right_cm:6.2f} B={bottom_cm:6.2f}{img_flag}{text_preview}")
    
    return all_shapes

# Analyze reference PPT (first few content slides)
ref_path = 'projects/GeCoM-Net_ppt169_20240616/templates/基于几何对应的多模态学习用于眼科图像分析 - 组会汇报.pptx'
ref_prs = Presentation(ref_path)

print("### REFERENCE PPT ALIGNMENT ###")
for i in [0, 2, 3, 4, 15]:  # Cover, content slides, last slide
    if i < len(ref_prs.slides):
        analyze_slide(ref_prs, i, 'REF')

# Analyze target PPT
tgt_path = 'projects/GeCoM-Net_ppt169_20240616/exports/GeCoM-Net_修改版.pptx'
tgt_prs = Presentation(tgt_path)

print("\n\n### TARGET PPT ALIGNMENT ###")
for i in range(min(19, len(tgt_prs.slides))):
    analyze_slide(tgt_prs, i, 'TGT')
