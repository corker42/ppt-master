import sys, os, traceback
SCRIPTS = os.path.join('G:', os.sep, 'ppt-master-main', 'skills', 'ppt-master', 'scripts')
sys.path.insert(0, SCRIPTS)
try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    pptx_path = os.path.join('G:', os.sep, 'ppt-master-main', 'projects', 'GeCoM-Net_ppt169_20240616', 'exports', 'GeCoM-Net_修改版.pptx')
    out_path = os.path.join('G:', os.sep, 'ppt-master-main', 'projects', 'GeCoM-Net_ppt169_20240616', 'exports', 'GeCoM-Net_修改版_v2.pptx')
    prs = Presentation(pptx_path)
    DG = RGBColor(0x3F, 0x3F, 0x3F)
    RA = RGBColor(0xC4, 0x1E, 0x3A)
    WH = RGBColor(0xFF, 0xFF, 0xFF)
    LG = RGBColor(0x61, 0x61, 0x61)
    LP = RGBColor(0xFD, 0xEC, 0xEF)
    GR = RGBColor(0x2E, 0x7D, 0x32)
    def atb(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        return tb.text_frame
    def sr(para, text, fs=13, b=False, i=False, c=DG):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(fs)
        r.font.bold = b
        r.font.italic = i
        r.font.color.rgb = c
        return r
    s9 = prs.slides[8]
    tf = atb(s9, Emu(914400), Emu(6950000), Emu(5000000), Emu(350000))
    p = tf.paragraphs[0]
    sr(p, 'Focal Loss ', 16, True, False, RA)
    sr(p, '公式详解', 16, True, False, DG)
    bg = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(810239), Emu(7300000), Emu(14600000), Emu(500000))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LP
    bg.line.fill.background()
    tf2 = bg.text_frame
    tf2.word_wrap = True
    tf2.margin_top = Emu(60000)
    tf2.margin_left = Emu(150000)
    pf = tf2.paragraphs[0]
    pf.alignment = PP_ALIGN.CENTER
    sr(pf, 'FL(p', 15, True, False, DG)
    sr(pf, 't', 12, True, True, RA)
    sr(pf, ') = ', 15, True, False, DG)
    sr(pf, '-α', 15, True, False, RA)
    sr(pf, ' · (1 - p', 15, True, False, DG)
    sr(pf, 't', 12, True, True, RA)
    sr(pf, ')', 15, True, False, DG)
    sr(pf, 'ᵞ', 18, True, False, RA)
    sr(pf, ' · log(p', 15, True, False, DG)
    sr(pf, 't', 12, True, True, RA)
    sr(pf, ')', 15, True, False, DG)
    tp = atb(s9, Emu(914400), Emu(7850000), Emu(14400000), Emu(500000))
    p1 = tp.paragraphs[0]
    sr(p1, 'α', 13, True, True, RA)
    sr(p1, ' = 类别平衡权重，用于平衡正负样本；  ', 12, False, False, LG)
    sr(p1, 'γ', 13, True, True, RA)
    sr(p1, ' = 聚焦参数（本项目 γ=2.0），控制难样本的权重放大程度；  ', 12, False, False, LG)
    sr(p1, 'p', 13, True, True, RA)
    sr(p1, 't', 10, True, True, RA)
    sr(p1, ' = 模型对真实类别的预测概率', 12, False, False, LG)
    p2 = tp.add_paragraph()
    sr(p2, '当 γ=0 时，Focal Loss 退化为标准交叉熵；γ 越大，对易分类样本的惩罚越强，使模型更专注于难分类样本', 11, False, False, LG)
    print('Slide 9 done')
    s12 = prs.slides[11]
    rm = [s for s in s12.shapes if s.name in ['图片 16', '图片 19', '图片 21']]
    sp = s12.shapes._spTree
    for s in rm:
        sp.remove(s._element)
        print('Removed: ' + s.name)
    tk = atb(s12, Emu(500000), Emu(1350000), Emu(7500000), Emu(400000))
    pk = tk.paragraphs[0]
    sr(pk, '评价指标 ', 18, True, False, RA)
    sr(pk, "Cohen\u2019s Kappa Coefficient (κ)", 18, True, False, DG)
    td = atb(s12, Emu(500000), Emu(1780000), Emu(7500000), Emu(500000))
    pd = td.paragraphs[0]
    sr(pd, "Cohen\u2019s kappa 系数是多分类任务的标准评价指标，取值范围 [-1, 1]，值越大表示一致性越好。", 12, False, False, LG)
    kb = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(500000), Emu(2300000), Emu(7200000), Emu(380000))
    kb.fill.solid()
    kb.fill.fore_color.rgb = LP
    kb.line.fill.background()
    kf = kb.text_frame
    kf.word_wrap = True
    kf.margin_top = Emu(50000)
    kf.margin_left = Emu(150000)
    pkf = kf.paragraphs[0]
    pkf.alignment = PP_ALIGN.CENTER
    sr(pkf, 'κ = (p', 16, True, False, DG)
    sr(pkf, '0', 11, True, False, DG)
    sr(pkf, ' - p', 16, True, False, DG)
    sr(pkf, 'e', 11, True, False, DG)
    sr(pkf, ') / (1 - p', 16, True, False, DG)
    sr(pkf, 'e', 11, True, False, DG)
    sr(pkf, ')', 16, True, False, DG)
    tdf = atb(s12, Emu(500000), Emu(2750000), Emu(7500000), Emu(900000))
    pd1 = tdf.paragraphs[0]
    sr(pd1, 'p', 13, True, True, RA)
    sr(pd1, '0', 10, True, True, RA)
    sr(pd1, ' = 混淆矩阵对角线元素之和 / 矩阵所有元素之和 = Accuracy', 12, False, False, DG)
    pd2 = tdf.add_paragraph()
    sr(pd2, 'p', 13, True, True, RA)
    sr(pd2, 'e', 10, True, True, RA)
    sr(pd2, ' = Σ(第i行元素之和 × 第i列元素之和) / (矩阵所有元素之和)²', 12, False, False, DG)
    pd3 = tdf.add_paragraph()
    sr(pd3, '', 6)
    pd4 = tdf.add_paragraph()
    sr(pd4, 'Score', 13, True, False, RA)
    sr(pd4, 'task1', 10, True, False, RA)
    sr(pd4, ' = κ × 10', 13, True, False, DG)
    tl1 = atb(s12, Emu(500000), Emu(3500000), Emu(7000000), Emu(280000))
    pl1 = tl1.paragraphs[0]
    sr(pl1, 'GAMMA 挑战赛提交记录（最佳 Score: 7.6138）', 13, True, False, RA)
    ts1 = s12.shapes.add_table(6, 4, Emu(500000), Emu(3800000), Emu(7500000), Emu(2200000))
    t1 = ts1.table
    for i, w in enumerate([Emu(1500000), Emu(1200000), Emu(1200000), Emu(3600000)]):
        t1.columns[i].width = w
    for i, h in enumerate(['Score', 'Kappa', '状态', '提交时间']):
        c = t1.cell(0, i)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WH
            p.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = DG
    d1 = [['7.6138','0.76138','已完成','2026-06-17 17:57'],['7.6138','0.76138','已完成','2026-06-15 17:54'],['7.35367','0.73537','已完成','2026-06-14 10:56'],['7.1991','0.71991','已完成','2026-06-18 08:58'],['6.79113','0.67911','已完成','2026-06-13 18:35']]
    for ri, rd in enumerate(d1):
        bgc = WH if ri % 2 == 0 else LP
        for ci, v in enumerate(rd):
            c = t1.cell(ri+1, ci)
            c.text = v
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = GR if ci == 2 else DG
                if ci == 2: p.font.bold = True
            c.fill.solid()
            c.fill.fore_color.rgb = bgc
    tl2 = atb(s12, Emu(8500000), Emu(3500000), Emu(6000000), Emu(280000))
    pl2 = tl2.paragraphs[0]
    sr(pl2, '早期提交记录（最佳 Kappa: 0.84232）', 13, True, False, RA)
    ts2 = s12.shapes.add_table(3, 4, Emu(8500000), Emu(3800000), Emu(7200000), Emu(1200000))
    t2 = ts2.table
    for i, w in enumerate([Emu(1500000), Emu(1200000), Emu(1200000), Emu(3300000)]):
        t2.columns[i].width = w
    for i, h in enumerate(['Score', 'Kappa', '状态', '提交时间']):
        c = t2.cell(0, i)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WH
            p.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = DG
    d2 = [['8.42317','0.84232','已完成','2026-06-11 12:30'],['7.15278','0.71528','已完成','2026-06-12 19:42']]
    for ri, rd in enumerate(d2):
        bgc = WH if ri % 2 == 0 else LP
        for ci, v in enumerate(rd):
            c = t2.cell(ri+1, ci)
            c.text = v
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = GR if ci == 2 else DG
                if ci == 2: p.font.bold = True
            c.fill.solid()
            c.fill.fore_color.rgb = bgc
    ts = atb(s12, Emu(8500000), Emu(5200000), Emu(7200000), Emu(600000))
    ps1 = ts.paragraphs[0]
    sr(ps1, '综合结果：', 13, True, False, DG)
    ps2 = ts.add_paragraph()
    sr(ps2, '• 最佳 Kappa: 0.84232 (Score: 8.42)', 12, False, False, LG)
    ps3 = ts.add_paragraph()
    sr(ps3, '• 最稳定 Score: 7.6138 (Kappa: 0.76)', 12, False, False, LG)
    ps4 = ts.add_paragraph()
    sr(ps4, '• 共提交 7 次，均已完成', 12, False, False, LG)
    print('Slide 12 done')
    prs.save(out_path)
    print('Saved: ' + out_path)
    prs2 = Presentation(out_path)
    print('Slides: %d' % len(prs2.slides))
    s12v = prs2.slides[11]
    pics = sum(1 for s in s12v.shapes if s.shape_type == 13 and s.name.startswith('图片'))
    print('Slide 12 pics remaining: %d' % pics)
    print('DONE')
except Exception:
    traceback.print_exc()
